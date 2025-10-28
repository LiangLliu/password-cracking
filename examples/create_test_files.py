#!/usr/bin/env python3
"""
Create password-protected test files for the password cracking tool.
Password: 92eo

Usage with UV:
    # First time - install dependencies:
    uv sync

    # Run the script:
    uv run python create_test_files.py
"""

import os
import sys
import subprocess
import shutil

PASSWORD = "92eo"

def create_test_pdf():
    """Create a password-protected PDF using PyPDF2."""
    try:
        from PyPDF2 import PdfWriter, PdfReader
        from reportlab.lib.pagesizes import letter
        from reportlab.pdfgen import canvas

        # First create a simple PDF
        temp_pdf = "temp_test.pdf"
        c = canvas.Canvas(temp_pdf, pagesize=letter)
        c.drawString(100, 750, "Password Cracking Test Document")
        c.drawString(100, 730, "This is a test PDF file.")
        c.drawString(100, 710, f"Password: {PASSWORD}")
        c.drawString(100, 690, "Created for testing the password cracking tool.")
        c.save()

        # Now encrypt it
        reader = PdfReader(temp_pdf)
        writer = PdfWriter()

        for page in reader.pages:
            writer.add_page(page)

        # Add password protection
        writer.encrypt(PASSWORD)

        # Write to current directory (examples/)
        output_path = "test.pdf"
        with open(output_path, "wb") as f:
            writer.write(f)

        # Clean up temp file
        os.remove(temp_pdf)
        print(f"✓ Created password-protected PDF: {output_path}")
        return output_path
    except ImportError as e:
        print(f"❌ Missing dependency for PDF: {e}")
        print("  Make sure to run with all dependencies (see usage at top of file)")
        return None
    except Exception as e:
        print(f"❌ Error creating PDF: {e}")
        return None

def create_test_docx():
    """Create a Word document and optionally encrypt with msoffice-crypt."""
    try:
        from docx import Document

        # Create a simple document
        doc = Document()
        doc.add_heading("Password Cracking Test Document", 0)
        doc.add_paragraph("This is a test Word document.")
        doc.add_paragraph(f"Password: {PASSWORD}")
        doc.add_paragraph("Created for testing the password cracking tool.")

        # Save without password first
        output_path = "test.docx"
        doc.save(output_path)
        print(f"✓ Created Word document: {output_path}")

        # Try to encrypt with msoffice-crypt if available
        if encrypt_office_file(output_path):
            print(f"✓ Encrypted Word document with password: {PASSWORD}")
        else:
            print("  Note: Manual password protection needed (msoffice-crypt not available)")

        return output_path
    except ImportError as e:
        print(f"❌ Missing dependency for Word: {e}")
        return None
    except Exception as e:
        print(f"❌ Error creating Word document: {e}")
        return None

def create_test_xlsx():
    """Create an Excel file and optionally encrypt with msoffice-crypt."""
    try:
        import openpyxl

        # Create a simple workbook
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Test Sheet"

        # Add some data
        ws['A1'] = "Password Cracking Test"
        ws['A2'] = "This is a test Excel file"
        ws['A3'] = f"Password: {PASSWORD}"
        ws['A4'] = "Created for testing the password cracking tool"

        # Add some sample data
        ws['A6'] = "Sample Data:"
        for i in range(7, 12):
            ws[f'A{i}'] = f"Row {i-6}"
            ws[f'B{i}'] = i * 100

        # Save without password first
        output_path = "test.xlsx"
        wb.save(output_path)
        print(f"✓ Created Excel file: {output_path}")

        # Try to encrypt with msoffice-crypt if available
        if encrypt_office_file(output_path):
            print(f"✓ Encrypted Excel file with password: {PASSWORD}")
        else:
            print("  Note: Manual password protection needed (msoffice-crypt not available)")

        return output_path
    except ImportError as e:
        print(f"❌ Missing dependency for Excel: {e}")
        return None
    except Exception as e:
        print(f"❌ Error creating Excel file: {e}")
        return None

def create_test_pptx():
    """Create a PowerPoint presentation and optionally encrypt with msoffice-crypt."""
    try:
        from pptx import Presentation

        # Create a simple presentation
        prs = Presentation()

        # Add a title slide
        slide_layout = prs.slide_layouts[0]  # Title slide
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "Password Cracking Test"
        subtitle.text = f"This is a test PowerPoint file\nPassword: {PASSWORD}"

        # Add a second slide with content
        slide_layout = prs.slide_layouts[1]  # Title and content
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = "Test Content"
        content.text = "This presentation was created for testing the password cracking tool.\n\n" \
                      f"The password is: {PASSWORD}\n\n" \
                      "This tests the tool's ability to crack PowerPoint files."

        # Save without password first
        output_path = "test.pptx"
        prs.save(output_path)
        print(f"✓ Created PowerPoint file: {output_path}")

        # Try to encrypt with msoffice-crypt if available
        if encrypt_office_file(output_path):
            print(f"✓ Encrypted PowerPoint file with password: {PASSWORD}")
        else:
            print("  Note: Manual password protection needed (msoffice-crypt not available)")

        return output_path
    except ImportError as e:
        print(f"❌ Missing dependency for PowerPoint: {e}")
        return None
    except Exception as e:
        print(f"❌ Error creating PowerPoint file: {e}")
        return None

def encrypt_office_file(filepath):
    """Encrypt an Office file using msoffcrypto-tool if available."""
    try:
        from msoffcrypto.format.ooxml import OOXMLFile

        # Create encrypted version
        temp_path = filepath + '.encrypted'

        with open(filepath, 'rb') as f:
            officefile = OOXMLFile(f)

            with open(temp_path, 'wb') as output:
                officefile.encrypt(PASSWORD, output)

        # Replace original with encrypted
        shutil.move(temp_path, filepath)
        return True

    except ImportError:
        # msoffcrypto-tool not available
        return False
    except Exception as e:
        print(f"  Encryption error: {e}")
        return False

def main():
    print("Password-Protected Test File Generator")
    print("=" * 40)
    print(f"Password for all files: {PASSWORD}")
    print("=" * 40)
    print()

    print("Creating test files...")
    print("-" * 40)

    # Create test files
    files_created = []

    pdf = create_test_pdf()
    if pdf:
        files_created.append(pdf)

    docx = create_test_docx()
    if docx:
        files_created.append(docx)

    xlsx = create_test_xlsx()
    if xlsx:
        files_created.append(xlsx)

    pptx = create_test_pptx()
    if pptx:
        files_created.append(pptx)

    # Summary
    print("\n" + "=" * 40)
    print("Summary:")
    print(f"Files created: {len(files_created)}")

    if files_created:
        print("\nTest files ready:")
        for f in files_created:
            print(f"  - {f}")

        print(f"\nAll files use password: {PASSWORD}")
        print("\nTo test with the password cracking tool:")
        print(f"../target/release/password-cracking -f {files_created[0]} dictionary -w ../wordlists/common-passwords.txt")

    # Note about manual encryption
    office_files = [f for f in files_created if f.endswith(('.docx', '.xlsx', '.pptx'))]
    if office_files:
        try:
            import msoffcrypto
        except ImportError:
            print("\n" + "=" * 40)
            print("Note: Office files were created but not encrypted.")
            print("Encryption should work automatically if you ran 'uv sync'.")
            print("\nAlternatively, add passwords manually in Microsoft Office.")

if __name__ == "__main__":
    main()